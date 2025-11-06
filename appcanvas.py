#!/usr/bin/env python3
import os
import sys
import re
import logging
import requests
from datetime import datetime
from fastapi import FastAPI, HTTPException, Query, Header, BackgroundTasks
from pydantic import BaseModel
from typing import Dict, Optional
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import io

STUDY_FILE_TYPES = {'pdf', 'pptx', 'ppt', 'docx', 'doc'}
VIDEO_FILE_TYPES = {'mp4', 'mov', 'avi', 'mkv', 'wmv', 'flv', 'webm'}

def extract_text_from_file(file_url: str, file_type: str) -> str:
    """Download and extract text from study materials."""
    try:
        response = requests.get(file_url, headers=headers, timeout=60)
        file_bytes = io.BytesIO(response.content)
        
        if file_type == 'pdf':
            reader = PdfReader(file_bytes)
            text = "\n".join(page.extract_text() for page in reader.pages)
            return text
        
        elif file_type in ['pptx', 'ppt']:
            prs = Presentation(file_bytes)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        elif file_type in ['docx', 'doc']:
            doc = Document(file_bytes)
            text = "\n".join(para.text for para in doc.paragraphs)
            return text
        
        return ""
    except Exception as e:
        logging.error(f"Could not extract text from {file_url}: {e}")
        return ""

def dump_course_smart(course_id: int) -> dict:
    """Smart sync - extracts text from PDFs/slides, skips videos."""
    course = _get(f"{CANVAS}courses/{course_id}")
    
    modules_url = f"{CANVAS}courses/{course_id}/modules"
    modules = _get(modules_url, include=["items"])
    
    result = {
        "course_id": course_id,
        "course_name": course.get("name"),
        "modules": []
    }
    
    for module in modules:
        module_data = {
            "module_id": module["id"],
            "module_name": module["name"],
            "items": []
        }
        
        items_url = f"{CANVAS}courses/{course_id}/modules/{module['id']}/items"
        items = _get(items_url)
        
        for item in items:
            item_data = {
                "item_id": item["id"],
                "title": item["title"],
                "type": item["type"],
                "url": item.get("url"),
                "html_url": item.get("html_url"),
                "content": None
            }
            
            if item["type"] == "Page" and "url" in item:
                page = _get(item["url"])
                item_data["content"] = page.get("body")
            
            elif item["type"] == "Assignment" and "content_id" in item:
                assignment = _get(f"{CANVAS}courses/{course_id}/assignments/{item['content_id']}")
                item_data["content"] = assignment.get("description")
            
            elif item["type"] == "File" and "content_id" in item:
                file_info = _get(f"{CANVAS}files/{item['content_id']}")
                filename = file_info.get("filename", "")
                file_extension = filename.split(".")[-1].lower()
                
                item_data["file_url"] = file_info.get("url")
                item_data["filename"] = filename
                
                if file_extension in STUDY_FILE_TYPES:
                    logging.info(f"Extracting text from {filename}")
                    item_data["content"] = extract_text_from_file(file_info.get("url"), file_extension)
                elif file_extension in VIDEO_FILE_TYPES:
                    logging.info(f"Skipping video: {filename}")
                    item_data["content"] = f"[Video: {filename}]"
            
            module_data["items"].append(item_data)
        
        result["modules"].append(module_data)
    
    return result

# Make the canvas_dump repo importable
# Adjust this path if your repo lives elsewhere.
sys.path.append("/Users/maxwellmoroz/canvastest/canvas_dump")

# These come from the repo
from canvas_dump_lib import headers, CANVAS, read_folder
from canvas_dump_utils import mk_and_cd  # we will provide a clean one below

logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")
app = FastAPI(title="Canvas Dump API")

# API Key Authentication
API_KEY = os.getenv("SUBJECTFOCUS_API_KEY")

# Simple in-memory job tracking
job_status: Dict[str, dict] = {}

COURSE_FOLDERS_URL = f"{CANVAS}courses/{{course_id}}/folders/root"

def safe_dirname(s: str) -> str:
    """Sanitize strings for directory names."""
    s = re.sub(r'[\\/<>:"|?*]+', "_", s)
    s = s.replace("\0", "_")
    s = re.sub(r"\s+", " ", s).strip()
    return s[:120] if len(s) > 120 else s

def require_api_key(x_api_key: Optional[str] = Header(None)):
    """Validate API key from header."""
    if not API_KEY:
        # If no API key is set, allow all requests (dev mode)
        logging.warning("No SUBJECTFOCUS_API_KEY set - authentication disabled!")
        return
    if x_api_key != API_KEY:
        raise HTTPException(status_code=401, detail="Unauthorized - invalid API key")

def _get(url: str, **params):
    """Requests GET with Canvas headers and good errors."""
    r = requests.get(url, headers=headers, params=params or None, timeout=30)
    if not r.ok:
        try:
            detail = r.json()
        except Exception:
            detail = r.text
        logging.error("GET %s -> %s %s", url, r.status_code, detail)
        raise HTTPException(status_code=r.status_code, detail=detail)
    return r.json()

def dump_course(course_id: int) -> dict:
    """
    Main dump function - downloads all files from a Canvas course.
    Returns a dict with status information.
    """
    import json

    # Save original directory
    original_dir = os.getcwd()

    try:
        # base output dir
        mk_and_cd("output")

        # optional: lookup course name for folder
        course = _get(f"{CANVAS}courses/{course_id}", include=["name", "course_code"])
        name_raw = course.get("name") or course.get("course_code") or f"course_{course_id}"
        name = safe_dirname(name_raw)
        logging.info('Processing "%s" (%s)', name, course_id)

        mk_and_cd(name)

        # get course root folder then recurse via read_folder
        root_folder = _get(COURSE_FOLDERS_URL.format(course_id=course_id))

        # write folder metadata at root for convenience
        with open("course_meta.json", "w") as f:
            f.write(json.dumps(course, indent=2))

        read_folder(root_folder)

        # back to output, then back to original
        os.chdir(original_dir)

        return {
            "course_id": course_id,
            "course_name": name,
            "status": "completed",
            "output_dir": os.path.join(os.getcwd(), "output", name)
        }
    except Exception as e:
        os.chdir(original_dir)
        logging.exception(f"Error dumping course {course_id}")
        return {
            "course_id": course_id,
            "status": "failed",
            "error": str(e)
        }

def run_dump_background(job_id: str, course_id: int):
    """Background task wrapper that updates job status."""
    job_status[job_id] = {
        "job_id": job_id,
        "course_id": course_id,
        "status": "running",
        "started_at": datetime.utcnow().isoformat(),
    }

    try:
        result = dump_course(course_id)
        job_status[job_id].update({
            "status": result.get("status", "completed"),
            "completed_at": datetime.utcnow().isoformat(),
            "result": result
        })
    except Exception as e:
        logging.exception(f"Background job {job_id} failed")
        job_status[job_id].update({
            "status": "failed",
            "completed_at": datetime.utcnow().isoformat(),
            "error": str(e)
        })

class DumpBody(BaseModel):
    course_id: int

@app.get("/health")
def health():
    """Health check - verifies Canvas API connection."""
    try:
        me = _get(f"{CANVAS}users/self/profile")
        return {
            "ok": True,
            "user": me.get("name"),
            "domain": CANVAS,
            "auth_required": bool(API_KEY)
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}

@app.get("/dump")
def dump_get_blocking(
    course_id: int = Query(..., gt=0),
    auth=require_api_key
):
    """
    Blocking dump - waits until finished, then returns result.
    Use this for small courses or when you need immediate confirmation.
    """
    try:
        return dump_course(course_id)
    except HTTPException:
        raise
    except Exception as e:
        logging.exception("Unhandled error for %s", course_id)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/dump")
def dump_post_blocking(body: DumpBody, auth=require_api_key):
    """Blocking dump via POST."""
    return dump_get_blocking(body.course_id)

@app.get("/dump/async")
def dump_get_async(
    course_id: int = Query(..., gt=0),
    background_tasks: BackgroundTasks = None,
    auth=require_api_key
):
    """
    Non-blocking dump - returns immediately with job_id.
    Use GET /status/{job_id} to check progress.
    Recommended for large courses.
    """
    import uuid
    job_id = f"job_{course_id}_{uuid.uuid4().hex[:8]}"

    background_tasks.add_task(run_dump_background, job_id, course_id)

    return {
        "accepted": True,
        "job_id": job_id,
        "course_id": course_id,
        "message": f"Job started. Check status at GET /status/{job_id}"
    }

@app.post("/dump/async")
def dump_post_async(
    body: DumpBody,
    background_tasks: BackgroundTasks,
    auth=require_api_key
):
    """Non-blocking dump via POST."""
    return dump_get_async(body.course_id, background_tasks)

@app.get("/status/{job_id}")
def get_status(job_id: str, auth=require_api_key):
    """Check the status of a background job."""
    if job_id not in job_status:
        raise HTTPException(status_code=404, detail=f"Job {job_id} not found")
    return job_status[job_id]

@app.get("/jobs")
def list_jobs(auth=require_api_key):
    """List all jobs (in-memory, lost on restart)."""
    return {"jobs": list(job_status.values())}
